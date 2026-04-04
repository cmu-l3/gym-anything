#!/usr/bin/env python3
"""
Verifier for Replace Presentation Fonts task.
Parses the PPTX file to verify font changes and content preservation.
"""

import json
import os
import sys
import tempfile
import logging
import shutil

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Try to import python-pptx, installing if necessary (for host environment)
try:
    from pptx import Presentation
except ImportError:
    logger.warning("python-pptx not found, attempting to install...")
    import subprocess
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
        from pptx import Presentation
    except Exception as e:
        logger.error(f"Failed to install python-pptx: {e}")
        # Proceeding will likely fail, but we'll try/catch below


def verify_replace_fonts(traj, env_info, task_info):
    """
    Verify that Courier New has been replaced with Liberation Sans
    while preserving text content.
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    metadata = task_info.get('metadata', {})
    forbidden_font = metadata.get('forbidden_font', 'Courier New').lower()
    required_font = metadata.get('required_font', 'Liberation Sans').lower()

    # 1. Retrieve task result JSON
    result_file = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", result_file.name)
        with open(result_file.name, 'r') as f:
            task_result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to load task results: {e}"}
    finally:
        if os.path.exists(result_file.name):
            os.unlink(result_file.name)

    # 2. Check basic file status
    if not task_result.get('file_exists'):
        return {
            "passed": False,
            "score": 0,
            "feedback": "Target PPTX file not found."
        }
    
    if not task_result.get('file_modified'):
        # Anti-gaming: file must be saved
        return {
            "passed": False,
            "score": 0,
            "feedback": "File was not saved/modified after task started."
        }

    # 3. Retrieve and Parse PPTX file
    pptx_temp = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    try:
        copy_from_env(task_result['file_path'], pptx_temp.name)
        prs = Presentation(pptx_temp.name)
    except Exception as e:
        return {
            "passed": False,
            "score": 10,
            "feedback": f"File exists but is corrupt or invalid: {e}"
        }

    # 4. Analyze Fonts and Content
    score = 15  # Base score for valid file + modified
    feedback_parts = ["File saved and valid"]
    
    total_runs = 0
    forbidden_count = 0
    required_count = 0
    inherited_count = 0  # Fonts set to None (inherit from theme)
    text_content = []

    try:
        # Traverse all slides and shapes
        slide_count = len(prs.slides)
        if slide_count == 4:
            score += 10
            feedback_parts.append("Slide count correct (4)")
        else:
            feedback_parts.append(f"Incorrect slide count: {slide_count}/4")

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        # Collect text for preservation check
                        text = paragraph.text.strip()
                        if text:
                            text_content.append(text)
                        
                        for run in paragraph.runs:
                            # Skip empty runs
                            if not run.text.strip():
                                continue
                                
                            total_runs += 1
                            font_name = run.font.name
                            
                            if font_name:
                                f_lower = font_name.lower()
                                if forbidden_font in f_lower:
                                    forbidden_count += 1
                                elif required_font in f_lower:
                                    required_count += 1
                            else:
                                inherited_count += 1

        # Evaluate Fonts (50 pts total for font logic)
        if total_runs == 0:
            return {"passed": False, "score": score, "feedback": "Presentation appears empty (no text found)."}

        # Check Forbidden Font (Critical)
        if forbidden_count == 0:
            score += 30
            feedback_parts.append(f"✅ No '{forbidden_font}' found")
        else:
            feedback_parts.append(f"❌ '{forbidden_font}' still present in {forbidden_count} places")

        # Check Required Font
        # Note: If user updated Master Slide/Theme properly, runs might return None (inherited).
        # We accept explicit "Liberation Sans" OR "None" (implying theme inheritance, provided Courier is gone)
        # However, purely "None" might just mean "Default", so we prefer explicit or strong evidence.
        # Given the task description asks to "Replace instances", explicit assignment is most likely.
        
        # We give credit if:
        # 1. Mostly Liberation Sans
        # 2. OR Mostly inherited AND Courier is gone (benefit of doubt that theme was updated)
        
        explicit_ratio = required_count / total_runs
        inherited_ratio = inherited_count / total_runs
        
        if explicit_ratio > 0.8:
            score += 25
            feedback_parts.append(f"✅ '{required_font}' applied consistently")
        elif (explicit_ratio + inherited_ratio) > 0.9 and forbidden_count == 0:
            score += 20
            feedback_parts.append(f"✅ Fonts updated (likely via theme/inheritance)")
        elif explicit_ratio > 0.5:
            score += 10
            feedback_parts.append(f"⚠️ Partial '{required_font}' application ({required_count}/{total_runs})")
        else:
            feedback_parts.append(f"❌ '{required_font}' not significantly detected")

        # Evaluate Content Preservation (20 pts)
        # Check for key phrases
        key_phrases = [
            "Q3 Department Meeting",
            "Campaign Performance",
            "Upcoming Initiatives",
            "Action Items",
            "Website traffic increased"
        ]
        
        full_text_blob = " ".join(text_content).lower()
        phrases_found = sum(1 for phrase in key_phrases if phrase.lower() in full_text_blob)
        
        if phrases_found == len(key_phrases):
            score += 20
            feedback_parts.append("✅ Content preserved")
        elif phrases_found >= len(key_phrases) - 1:
            score += 10
            feedback_parts.append("⚠️ Most content preserved")
        else:
            feedback_parts.append("❌ Significant content loss")

    except Exception as e:
        logger.error(f"Error parsing PPTX: {e}")
        feedback_parts.append(f"Error during font analysis: {str(e)}")

    # Determine pass/fail
    # Must have removed forbidden font AND have a decent score
    passed = (forbidden_count == 0) and (score >= 70)

    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }