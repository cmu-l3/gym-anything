#!/usr/bin/env python3
"""
Verifier for Edit Slide Master Branding task.
Verifies ODP/PPTX structure to ensure master slides were edited correctly.
"""

import json
import os
import sys
import tempfile
import logging
import shutil
from pathlib import Path

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Attempt to import ODF/PPTX libraries (these are available in the env, 
# but for the verifier running on host, we might need to handle imports carefully
# or rely on what's available in the verifier's environment)
try:
    from odf import opendocument, style, draw, text as odftext
    ODF_AVAILABLE = True
except ImportError:
    ODF_AVAILABLE = False
    logger.warning("odfpy not available on verifier host")

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx not available on verifier host")


def verify_slide_master_branding(traj, env_info, task_info):
    """
    Verify that the user edited the slide master to add branding.
    
    Criteria:
    1. File exists and is valid (10 pts)
    2. Slide count (4) preserved (15 pts)
    3. Content preserved (titles match) (15 pts)
    4. Master page has new shapes (header bar) (20 pts)
    5. Master page has specific text ("Community Health Alliance") (25 pts)
    6. VLM Trajectory Check (15 pts)
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    metadata = task_info.get('metadata', {})
    expected_org_name = metadata.get('expected_org_name', "Community Health Alliance")
    
    score = 0
    max_score = 100
    feedback_parts = []
    
    # =========================================================
    # 1. Retrieve Result JSON and File
    # =========================================================
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.dat')
    
    try:
        copy_from_env("/tmp/task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result_data = json.load(f)
            
        output_exists = result_data.get("output_exists", False)
        file_format = result_data.get("file_format", "")
        
        if output_exists:
            try:
                copy_from_env("/tmp/submission_file.dat", temp_file.name)
                file_available = True
            except Exception as e:
                logger.error(f"Failed to copy submission file: {e}")
                file_available = False
        else:
            file_available = False

    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to retrieve task results: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    # =========================================================
    # 2. Evaluate File Existence (10 pts)
    # =========================================================
    if output_exists and file_available:
        score += 10
        feedback_parts.append(f"✅ File saved successfully ({file_format})")
    else:
        return {
            "passed": False, 
            "score": 0, 
            "feedback": "❌ No saved presentation file found. Did you save the file?"
        }

    # =========================================================
    # 3. Parse File and Analyze Content
    # =========================================================
    slide_count = 0
    titles_found = 0
    master_shapes_found = False
    org_name_on_master = False
    
    try:
        if file_format == 'odp' and ODF_AVAILABLE:
            doc = opendocument.load(temp_file.name)
            
            # Check slide count
            slides = doc.getElementsByType(draw.Page)
            slide_count = len(slides)
            
            # Check content preservation
            all_text_content = []
            for slide in slides:
                slide_text = []
                for t in slide.getElementsByType(odftext.P):
                    slide_text.append(str(t))
                all_text_content.append(" ".join(slide_text).lower())
                
            expected_titles = metadata.get("original_titles", [])
            for title in expected_titles:
                if any(title.lower() in slide_txt for slide_txt in all_text_content):
                    titles_found += 1
            
            # Check Master Page
            # In ODP, master pages are in doc.masterstyles
            # We look for shapes (rect, customshape) and text "Community Health Alliance"
            # specifically within master-page elements
            for mp in doc.masterstyles.childNodes:
                # Check for shapes (header bar)
                shapes = []
                shapes.extend(mp.getElementsByType(draw.Rect))
                shapes.extend(mp.getElementsByType(draw.CustomShape))
                # Default masters usually have just frames, so finding Rect/CustomShape is good
                if len(shapes) > 0:
                    master_shapes_found = True
                
                # Check for text in master
                mp_text_nodes = mp.getElementsByType(odftext.P)
                for node in mp_text_nodes:
                    if expected_org_name.lower() in str(node).lower():
                        org_name_on_master = True
            
        elif (file_format == 'pptx' or file_format == '') and PPTX_AVAILABLE:
            # Fallback for PPTX or if ODF not available but python-pptx is
            prs = Presentation(temp_file.name)
            slide_count = len(prs.slides)
            
            expected_titles = metadata.get("original_titles", [])
            for slide in prs.slides:
                txt = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        txt += shape.text + " "
                
                for title in expected_titles:
                    if title.lower() in txt.lower():
                        titles_found += 1
                        # Remove to avoid double counting if titles repeat
                        # (though they don't in this task)
                        pass 
            
            # Normalize titles count (simple check)
            titles_found = min(titles_found, 4)

            # Check Master
            for master in prs.slide_masters:
                # Check for explicit shapes on master
                # A standard blank master has placeholders. Added branding is usually a shape.
                # Heuristic: if shape count > initial baseline (which is usually small)
                # But safer: check if any shape contains the specific text
                for shape in master.shapes:
                    if hasattr(shape, "text") and expected_org_name.lower() in shape.text.lower():
                        org_name_on_master = True
                        master_shapes_found = True # If text is there, a shape is there
                    
                    # Also check for colored rectangles if text not found yet
                    if not master_shapes_found and shape.shape_type == 1: # MSO_SHAPE.RECTANGLE is auto_shape
                        master_shapes_found = True
                        
            # Also check slide layouts, users sometimes edit layout instead of master
            if not org_name_on_master:
                for layout in prs.slide_layouts:
                    for shape in layout.shapes:
                        if hasattr(shape, "text") and expected_org_name.lower() in shape.text.lower():
                            org_name_on_master = True
                            master_shapes_found = True

        else:
            feedback_parts.append(f"⚠️ Could not fully parse {file_format} file (missing libraries on verifier)")
            # Fallback: simple grep on the file content if it's XML based/unzipped, 
            # but binary files won't work. We'll trust file existence for partial points.
            pass

    except Exception as e:
        logger.error(f"Error parsing presentation: {e}")
        feedback_parts.append(f"⚠️ Error parsing file content: {e}")

    # Clean up temp file
    if os.path.exists(temp_file.name):
        os.unlink(temp_file.name)

    # =========================================================
    # 4. Score Calculation based on Analysis
    # =========================================================
    
    # Slide Count (15 pts)
    if slide_count == 4:
        score += 15
        feedback_parts.append("✅ Slide count preserved (4 slides)")
    else:
        feedback_parts.append(f"❌ Slide count changed (found {slide_count}, expected 4)")
        
    # Content Preservation (15 pts)
    if titles_found >= 3:
        score += 15
        feedback_parts.append("✅ Original content preserved")
    else:
        feedback_parts.append("⚠️ Some original content seems missing")

    # Master Shapes (20 pts)
    if master_shapes_found:
        score += 20
        feedback_parts.append("✅ New shapes detected on master/layout")
    else:
        feedback_parts.append("❌ No new shapes detected on master")

    # Org Name on Master (25 pts)
    if org_name_on_master:
        score += 25
        feedback_parts.append(f"✅ Organization name '{expected_org_name}' found on master")
    else:
        feedback_parts.append(f"❌ '{expected_org_name}' not found on master slide")

    # =========================================================
    # 5. VLM Trajectory Verification (15 pts)
    # =========================================================
    # Verify the workflow: "View > Master Slide" -> "Close Master View"
    from gym_anything.vlm import sample_trajectory_frames, query_vlm
    
    vlm_score = 0
    try:
        frames = sample_trajectory_frames(traj, n=4)
        if frames:
            prompt = """
            You are verifying a LibreOffice Impress task. 
            The user should have:
            1. Opened "Master Slide" view (look for "Close Master View" button in toolbar).
            2. Added a blue rectangle/header at the top of the slide.
            3. Added text "Community Health Alliance".
            
            Look at these frames. Do you see evidence of editing the Slide Master?
            - Is there a "Close Master View" button visible in any frame?
            - Do you see a blue header being drawn or appearing?
            - Is the text "Community Health Alliance" visible?
            
            Return JSON: {"master_view_entered": bool, "branding_added": bool}
            """
            
            vlm_res = query_vlm(prompt=prompt, images=frames)
            
            if vlm_res and vlm_res.get("success"):
                parsed = vlm_res.get("parsed", {})
                if parsed.get("master_view_entered"):
                    vlm_score += 10
                    feedback_parts.append("✅ VLM confirmed Master View usage")
                else:
                    feedback_parts.append("⚠️ VLM did not see Master View opened")
                    
                if parsed.get("branding_added"):
                    vlm_score += 5
                    feedback_parts.append("✅ VLM confirmed branding elements")
            else:
                # If VLM fails/unavailable, we grant points if programmatic check passed hard
                if org_name_on_master:
                    vlm_score = 15
                    feedback_parts.append("✅ VLM skipped (programmatic check strong)")
        else:
             feedback_parts.append("⚠️ No trajectory frames available for VLM")
             
    except Exception as e:
        logger.warning(f"VLM check failed: {e}")
        # Fallback if VLM errors
        if org_name_on_master:
            vlm_score = 15

    score += vlm_score

    # Final Pass/Fail
    passed = score >= 60 and org_name_on_master
    
    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }