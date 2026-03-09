#!/usr/bin/env python3
"""
Verifier for format_slide_titles task.
Checks standard formatting (Font, Size, Color, Bold) on 5 slides.
"""

import json
import os
import sys
import tempfile
import logging
import shutil
from typing import Dict, Any, Tuple

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

def verify_format_slide_titles(traj, env_info, task_info):
    """
    Verify that all 5 slide titles have been reformatted correctly.
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    # Get expected values
    metadata = task_info.get('metadata', {})
    expected_titles = metadata.get('expected_titles', [])
    target_style = metadata.get('target_style', {})
    tolerances = metadata.get('tolerances', {})

    target_font = target_style.get('font_family', 'Liberation Sans')
    target_size = target_style.get('font_size_pt', 36)
    target_color_hex = target_style.get('font_color_hex', '#003366')
    target_bold = target_style.get('font_bold', True)
    
    tol_size = tolerances.get('font_size_pt', 1.5)
    tol_color = tolerances.get('color_rgb', 20)

    # 1. Retrieve Result JSON
    temp_json = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_json.name)
        with open(temp_json.name, 'r') as f:
            result_meta = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to load result metadata: {e}"}
    finally:
        if os.path.exists(temp_json.name):
            os.unlink(temp_json.name)

    if not result_meta.get('file_exists') or not result_meta.get('file_modified'):
        return {
            "passed": False, 
            "score": 0, 
            "feedback": "Presentation file was not saved or modified."
        }

    file_format = result_meta.get('file_format', '')
    result_path = result_meta.get('result_path', '')
    
    if not result_path or not file_format:
         return {
            "passed": False, 
            "score": 0, 
            "feedback": "No valid result file found in export."
        }

    # 2. Retrieve Presentation File
    temp_pres = tempfile.NamedTemporaryFile(delete=False, suffix=f'.{file_format}')
    temp_pres_path = temp_pres.name
    temp_pres.close()

    try:
        copy_from_env(result_path, temp_pres_path)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to retrieve presentation file: {e}"}

    # 3. Parse Presentation
    parsed_slides = []
    error_msg = None

    try:
        if file_format == 'pptx':
            parsed_slides = parse_pptx(temp_pres_path, expected_titles)
        elif file_format == 'odp':
            parsed_slides = parse_odp(temp_pres_path, expected_titles)
        else:
            error_msg = f"Unknown file format: {file_format}"
    except Exception as e:
        error_msg = f"Parsing error: {str(e)}"
    finally:
        if os.path.exists(temp_pres_path):
            os.unlink(temp_pres_path)

    if error_msg:
        return {"passed": False, "score": 0, "feedback": error_msg}

    if not parsed_slides:
        return {"passed": False, "score": 0, "feedback": "Could not extract slides from file."}

    # 4. Score Logic
    score = 0
    feedback_lines = []
    
    # Base points for valid file
    score += 10 
    feedback_lines.append("File saved and valid (+10)")

    # Score each slide (Max 13 points per slide: 4 font, 3 size, 3 color, 3 bold)
    slides_formatted_correctly = 0
    
    for i, slide in enumerate(parsed_slides):
        if i >= 5: break
        
        slide_score = 0
        checks = []
        
        # Check Font Family
        font_actual = slide.get('font_name', '')
        # Case insensitive partial match
        if target_font.lower() in font_actual.lower():
            slide_score += 4
            checks.append("Font OK")
        else:
            checks.append(f"Font mismatch ({font_actual})")

        # Check Font Size
        size_actual = slide.get('font_size', 0)
        if abs(size_actual - target_size) <= tol_size:
            slide_score += 3
            checks.append("Size OK")
        else:
            checks.append(f"Size mismatch ({size_actual}pt)")

        # Check Color
        color_actual = slide.get('font_color', (0,0,0)) # RGB tuple
        target_rgb = hex_to_rgb(target_color_hex)
        if color_matches(color_actual, target_rgb, tol_color):
            slide_score += 3
            checks.append("Color OK")
        else:
            checks.append(f"Color mismatch ({color_actual})")

        # Check Bold
        bold_actual = slide.get('font_bold', False)
        if bold_actual == target_bold:
            slide_score += 3
            checks.append("Bold OK")
        else:
            checks.append("Bold mismatch")

        score += slide_score
        
        if slide_score == 13:
            slides_formatted_correctly += 1
            feedback_lines.append(f"Slide {i+1}: Perfect ({', '.join(checks)})")
        else:
            feedback_lines.append(f"Slide {i+1}: {slide_score}/13 pts. Issues: {', '.join([c for c in checks if 'mismatch' in c])}")

    # Bonus: All Consistent (10 pts)
    if slides_formatted_correctly == 5:
        score += 10
        feedback_lines.append("Bonus: All slides consistent (+10)")

    # Content Preservation (15 pts)
    # Check if titles generally match expected (allows for minor whitespace diffs)
    titles_preserved = 0
    for i, slide in enumerate(parsed_slides):
        if i < len(expected_titles):
            if expected_titles[i] in slide.get('text', ''):
                titles_preserved += 1
    
    if titles_preserved >= 5:
        score += 15
        feedback_lines.append("Content preserved (+15)")
    else:
        feedback_lines.append(f"Content modified/lost ({titles_preserved}/5 preserved)")

    return {
        "passed": score >= 60,
        "score": score,
        "feedback": "\n".join(feedback_lines),
        "details": {
            "file_format": file_format,
            "slides_perfect": slides_formatted_correctly
        }
    }


def hex_to_rgb(hex_str):
    hex_str = hex_str.lstrip("#")
    return tuple(int(hex_str[i:i+2], 16) for i in (0, 2, 4))

def color_matches(c1, c2, tolerance):
    if not c1 or not c2: return False
    return all(abs(a - b) <= tolerance for a, b in zip(c1, c2))

def parse_pptx(filepath, expected_titles):
    """Parse PPTX file using python-pptx"""
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        prs = Presentation(filepath)
        slides_data = []
        
        for slide in prs.slides:
            # Find title shape
            title_shape = slide.shapes.title
            if not title_shape:
                # Fallback to first text shape
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        title_shape = shape
                        break
            
            data = {
                "text": "", "font_name": "", "font_size": 0, 
                "font_color": (0,0,0), "font_bold": False
            }
            
            if title_shape and title_shape.has_text_frame:
                data["text"] = title_shape.text
                # Analyze first run of first paragraph for formatting
                if title_shape.text_frame.paragraphs:
                    p = title_shape.text_frame.paragraphs[0]
                    if p.runs:
                        r = p.runs[0]
                        data["font_name"] = r.font.name or "Unknown"
                        if r.font.size:
                            data["font_size"] = r.font.size.pt
                        
                        # Color is tricky in PPTX
                        if r.font.color and r.font.color.type == 1: # RGB
                            data["font_color"] = tuple(r.font.color.rgb)
                        # We ignore theme colors for this strict task, or assume black if not set
                        
                        data["font_bold"] = bool(r.font.bold)
            
            slides_data.append(data)
        return slides_data
    except ImportError:
        return []

def parse_odp(filepath, expected_titles):
    """Parse ODP file using odfpy"""
    try:
        from odf import opendocument, draw, text
        from odf.namespaces import FONS
        doc = opendocument.load(filepath)
        slides_data = []
        
        # Build style map
        styles = {}
        for style in doc.automaticstyles.childNodes:
            name = style.getAttribute('name')
            props = {}
            for child in style.childNodes:
                if child.tagName == 'style:text-properties':
                    props['font_name'] = child.getAttribute('fontname') or child.getAttrNS(FONS, 'font-family')
                    props['font_size'] = child.getAttrNS(FONS, 'font-size')
                    props['font_color'] = child.getAttrNS(FONS, 'color')
                    props['font_weight'] = child.getAttrNS(FONS, 'font-weight')
            styles[name] = props

        for page in doc.getElementsByType(draw.Page):
            data = {
                "text": "", "font_name": "", "font_size": 0, 
                "font_color": (0,0,0), "font_bold": False
            }
            
            # Find title frame (usually the first frame)
            frames = page.getElementsByType(draw.Frame)
            if frames:
                # Extract text and style from first paragraph of first frame
                # This is a simplification; robust parsing is harder
                title_frame = frames[0]
                text_content = []
                
                # Check text content
                for p in title_frame.getElementsByType(text.P):
                    text_content.append(str(p))
                    
                    # Check styles on spans
                    spans = p.getElementsByType(text.Span)
                    if spans:
                        first_span = spans[0]
                        style_name = first_span.getAttribute('stylename')
                        if style_name in styles:
                            s = styles[style_name]
                            data['font_name'] = s.get('font_name', '')
                            
                            sz = s.get('font_size', '0pt')
                            try:
                                data['font_size'] = float(sz.replace('pt', ''))
                            except: pass
                            
                            c = s.get('font_color', '#000000')
                            data['font_color'] = hex_to_rgb(c)
                            
                            w = s.get('font_weight', 'normal')
                            data['font_bold'] = (w == 'bold')
                
                data["text"] = " ".join(text_content)
                
            slides_data.append(data)
            
        return slides_data
    except ImportError:
        return []