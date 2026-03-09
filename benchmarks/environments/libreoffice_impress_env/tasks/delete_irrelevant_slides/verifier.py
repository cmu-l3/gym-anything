#!/usr/bin/env python3
"""
Verifier for delete_irrelevant_slides task.
Checks that exactly 3 specific slides were deleted from the presentation.
"""

import json
import os
import sys
import tempfile
import logging
import hashlib

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import environment-specific utils if available, or fallback
try:
    sys.path.insert(0, '/workspace/utils')
    from impress_verification_utils import parse_odp_file, parse_pptx_file
except ImportError:
    # Minimal fallback for PPTX if utils not found in verifier context
    try:
        from pptx import Presentation
        def parse_pptx_file(filepath):
            try:
                prs = Presentation(filepath)
                slides_data = []
                for s in prs.slides:
                    text_elements = []
                    if s.shapes.title and s.shapes.title.text:
                        text_elements.append(s.shapes.title.text)
                    for shape in s.shapes:
                        if hasattr(shape, "text") and shape.text and shape != s.shapes.title:
                            text_elements.append(shape.text)
                    slides_data.append({'text_elements': text_elements})
                return {'slides': slides_data, 'slide_count': len(slides_data)}
            except Exception as e:
                return {'error': str(e)}
    except ImportError:
        pass

def verify_delete_slides(traj, env_info, task_info):
    """
    Verify that 3 specific slides were deleted and 7 remain.
    """
    copy_from_env = env_info.get('copy_from_env')
    if not copy_from_env:
        return {"passed": False, "score": 0, "feedback": "Copy function not available"}

    # Load result metadata
    temp_result = tempfile.NamedTemporaryFile(delete=False, suffix='.json')
    try:
        copy_from_env("/tmp/task_result.json", temp_result.name)
        with open(temp_result.name, 'r') as f:
            result = json.load(f)
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to load result JSON: {e}"}
    finally:
        if os.path.exists(temp_result.name):
            os.unlink(temp_result.name)

    # Check basic file existence
    if not result.get('file_exists', False):
        return {"passed": False, "score": 0, "feedback": "Presentation file not found"}

    final_path = result.get('final_file_path')
    if not final_path:
        return {"passed": False, "score": 0, "feedback": "No output file path specified"}

    # Copy the actual presentation file
    ext = os.path.splitext(final_path)[1].lower()
    temp_pres = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
    try:
        copy_from_env(final_path, temp_pres.name)
        
        # Parse based on extension
        parsed_data = {}
        if ext == '.pptx':
            try:
                # Use local fallback if needed, or imported util
                parsed_data = parse_pptx_file(temp_pres.name)
            except NameError:
                 return {"passed": False, "score": 0, "feedback": "python-pptx not available for verification"}
        elif ext == '.odp':
            try:
                # Assuming odfpy is available in verifier env as it is in container
                from impress_verification_utils import parse_odp_file
                parsed_data = parse_odp_file(temp_pres.name)
            except (ImportError, NameError):
                 return {"passed": False, "score": 0, "feedback": "ODP parser not available for verification"}
        else:
            return {"passed": False, "score": 0, "feedback": f"Unsupported file format: {ext}"}
            
    except Exception as e:
        return {"passed": False, "score": 0, "feedback": f"Failed to retrieve/parse file: {e}"}
    finally:
        if os.path.exists(temp_pres.name):
            os.unlink(temp_pres.name)

    if 'error' in parsed_data:
        return {"passed": False, "score": 0, "feedback": f"File parsing error: {parsed_data['error']}"}

    # Scoring Logic
    score = 0
    feedback_parts = []
    
    # 1. File Validity (10 pts)
    score += 10
    feedback_parts.append("File valid")

    # 2. Slide Count (20 pts)
    slide_count = parsed_data.get('slide_count', 0)
    expected_count = 7
    if slide_count == expected_count:
        score += 20
        feedback_parts.append("Slide count correct (7)")
    else:
        feedback_parts.append(f"Incorrect slide count: {slide_count} (expected 7)")

    # 3. Content Verification
    metadata = task_info.get('metadata', {})
    deleted_titles = metadata.get('deleted_titles', [])
    retained_titles = metadata.get('retained_titles', [])
    
    # Extract all text for checking
    all_text_blob = ""
    slide_titles_found = []
    
    slides = parsed_data.get('slides', [])
    for slide in slides:
        text_elems = slide.get('text_elements', [])
        slide_text = " ".join(text_elems)
        all_text_blob += slide_text + "\n"
        if text_elems:
            slide_titles_found.append(text_elems[0]) # Assuming first element is title

    # Check deleted titles (15 pts each = 45 pts)
    deleted_count = 0
    for title in deleted_titles:
        if title not in all_text_blob:
            score += 15
            deleted_count += 1
        else:
            feedback_parts.append(f"Failed to delete: '{title}'")
            
    if deleted_count == 3:
        feedback_parts.append("All target slides deleted")

    # Check retained titles (5 pts each for key ones, capped at 20 pts total to fit 100 scale)
    # Actually design said 5 pts each for 4 specific slides + others. 
    # Let's align with the totals: 10+20+45 = 75 so far. 
    # Remaining 25 points: 5 for modified, 20 for retained.
    
    retained_score = 0
    # Check 4 key retained slides
    key_retained = [
        "Business Impact Analysis Report", 
        "Critical Business Functions", 
        "Recovery Time Objectives", 
        "Action Items and Next Steps"
    ]
    
    found_retained = 0
    for title in key_retained:
        if title in all_text_blob:
            found_retained += 1
    
    retained_score = found_retained * 5
    score += retained_score
    if found_retained == 4:
        feedback_parts.append("Key slides retained")
    else:
        feedback_parts.append(f"Missing {4-found_retained} key slides")

    # 4. Anti-Gaming (File Modified) (5 pts)
    if result.get('file_modified', False):
        score += 5
    else:
        feedback_parts.append("File not modified (timestamp check failed)")

    passed = (score >= 75) and (slide_count == 7)

    return {
        "passed": passed,
        "score": score,
        "feedback": " | ".join(feedback_parts)
    }